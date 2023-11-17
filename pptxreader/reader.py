from pptx import Presentation
import glob
import time
import textwrap

for eachfile in glob.glob("*.pptx"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    x = ""
    start = time.monotonic()
    with open("file.txt", "w") as f:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    x = x + " " + shape.text
        x = textwrap.wrap(x, 1024)
        y = ""
        for i in x:
            y = y + "\n\n" + i
        print(y)
        f.write(y)
    print(f"End: {time.monotonic()-start}")
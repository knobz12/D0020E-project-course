from pptx import Presentation
import time
import textwrap


class PPTXToText():
    def __init__(self):
        pass

    def ConvertToText(self, file):
        extracted_text = ""
        PPTX = Presentation(file)

        for slide in PPTX.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text = extracted_text + " " + shape.text
        return extracted_text
    
    def ConvertImageToText(self, file):
        pass
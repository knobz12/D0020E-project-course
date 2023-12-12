import textwrap
import magic

import modules.files.docx3txt as docx3txt

import pytesseract
from PIL import Image
from bs4 import BeautifulSoup

import pptx
import io
import PyPDF2

import os
import pathlib
import sys

import glob
import json

files_dir = pathlib.Path("./bin/jsoneler/files")

class NotSupportedFiletype(Exception):
    def __init__(self, message = "Type not currently supported. We are working on a fix in later versions."):
        super().__init__(message)

# TODO(Johan) kanske göra så text_and_image_text_from_file_bytes inte returnerar filetype
class Chunkerizer:
    
    def make_chunk(text: str, chunk_size: int) -> list[str]:
        return textwrap.wrap(text, chunk_size)

    def text_and_image_text_from_file(filepath: str, image_txt_from_complex_types: bool) -> tuple[str, str, str] | None:
        filename = os.path.basename(filepath)
        with open(filepath, "rb") as f:
            return Chunkerizer.text_and_image_text_from_file_bytes(f.read(), image_txt_from_complex_types, filename)

    def text_and_image_text_from_file_bytes(buf: bytes, image_txt_from_complex_types: bool, filename: str = "") -> tuple[str, str, str] | None:        
        def parse_html(buf: bytes) -> str:
            soup = BeautifulSoup(buf, features="html.parser")
            return soup.get_text()



        def is_programming_lang(mime_type: str) -> bool:
            return mime_type == "text/x-c++" or mime_type == "text/x-python" or mime_type == "text/x-java"

        file_extenstion = filename[filename.rfind(".") + 1:] 


        extracted_text = ""
        extracted_image_text = ""
        try:
            mime_type = magic.from_buffer(buf[0:512], mime = True)
            #print(mime_type)
            filetype = mime_type[mime_type.find("/") + 1:]

            if mime_type.startswith("image"):
                img = Image.open(io.BytesIO(buf))
                extracted_text = pytesseract.image_to_string(img)
                extracted_image_text = extracted_text
            elif mime_type == "application/json":
                mime_type = "text/plain" 

            elif mime_type == "application/pdf": 

                io_buf = io.BytesIO(buf)
                pdf = PyPDF2.PdfReader(io_buf)
                
                image_list = []

                for page in pdf.pages:
                    extracted_text += page.extract_text()
                    if image_txt_from_complex_types:
                        for img in page.images:
                            if img not in image_list:
                                image_list.append(img)

                if image_txt_from_complex_types:
                    for img_obj in image_list:
                        io_buf = io.BytesIO(img_obj.data)
                        img = Image.open(io_buf)
                        extracted_image_text += pytesseract.image_to_string(img)
                
            elif mime_type == "application/x-empty":
                pass
            elif is_programming_lang(mime_type):
                extracted_text = bytes.decode(buf, "utf-8", errors = "ignore")
            elif mime_type == "text/plain":
                if file_extenstion == "html":
                    extracted_text = parse_html(buf)
                    filetype = "html"
                elif file_extenstion == "json":
                    extracted_text = bytes.decode(buf, "utf-8", errors = "ignore")
                    filetype = "json"
                elif file_extenstion == "jsonl":
                    extracted_text = bytes.decode(buf, "utf-8", errors = "ignore")
                    filetype = "jsonl"
                else:
                    extracted_text = bytes.decode(buf, "utf-8", errors = "ignore")
            elif mime_type == "text/html":
                extracted_text = parse_html(buf)
            elif mime_type == "text/xml":
                soup = BeautifulSoup(buf, 'xml')
                extracted_text = soup.get_text()
            elif mime_type == "application/zip" or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                if file_extenstion == "zip":
                    # TODO(Johan) support zip correctly
                    raise NotSupportedFiletype

                filetype = "ppt"

                io_buf = io.BytesIO(buf)
                powerpoint = pptx.Presentation(io_buf)

                for slide in powerpoint.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text = extracted_text + " " + shape.text
                        elif(image_txt_from_complex_types and hasattr(shape, "image")):
                            io_buf = io.BytesIO(shape.image.blob)
                            img = Image.open(io_buf)
                            if img.format == "WMF" or img.format == "EMF":
                                continue
                            extracted_image_text += pytesseract.image_to_string(img)
                            img.close()
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                filetype = "docx"
                io_buf = io.BytesIO(buf)
                (extracted_text, extracted_image_text) = docx3txt.process(io_buf)
            else:
                raise NotSupportedFiletype

                    
            return filetype, extracted_text, extracted_image_text
            
        except Exception as error:
            print(type(error).__name__, "-", error)
            return None
        
    def text_extraction_test():
        #path_list = ["./backend/tests\sample_files\courses\D7032E\evolutionCards.json"]
        #path_list = ["./backend/tests/sample_files/Test_txt/t.txt"]
        path_list = glob.glob("./backend/tests/**/*.*", recursive = True)
        try:
            for path in path_list: 
                print(f"###### testing {path}")
                Chunkerizer.text_and_image_text_from_file(path, True)
        except Exception as error:
            pass
        





def get_course_codes() -> list[str]:
    # path = files_folder.resolve()
    dirs = os.listdir(files_dir.resolve())
    folders: list[str] = []
    for folder in dirs:
        folders.append(folder)
    return folders

def get_course_files(code: str) -> list[str]:
    code_dir = pathlib.Path(files_dir, code)
    files = os.listdir(code_dir)
    return files

def get_file_hash(text: str) -> str:
    from hash_and_json.HashFunction import TextToHash
    (hash, _) =  TextToHash().ConvertToHash(text)
    return hash

def create_json_chunks(filehash: str, text_chunks: list[str], course: str) -> list[str]:
    import json
    jsons: list[str] = []
    for (idx, chunk) in enumerate(text_chunks):
        if len(chunk) < 64:
            continue
        obj = {
            "id":filehash,
            "chunk-id":idx,
            "course": course,
            "text": chunk,
        }
        result = json.dumps(obj,indent=0).replace("\n","")
        jsons.append(result)
    return jsons

def create_jsonls():
    courses = get_course_codes()
    print("Courses:",courses)
    jsonls: list[str] = []

    for course in courses[0:1]:
        code_dir = pathlib.Path(files_dir, course)
        files = get_course_files(course)
        # for file in files:
        for file in files:
            filepath = pathlib.Path(code_dir, file).resolve()
            print(filepath)

            result = Chunkerizer.text_and_image_text_from_file(str(filepath), False)

            if result == None:
                continue

            (_, text) = result

            chunks = Chunkerizer.make_chunk(text, 512)
            file_hash = get_file_hash(text)

            if file_hash == None:
                continue

            jsonl = create_json_chunks(file_hash, chunks, course)
            jsonls.extend(jsonl)
    return jsonls

if __name__ == "__main__":
    jsonls = create_jsonls()
    result_text = ""
    for jsonl in jsonls:
        result_text += jsonl + "\n"

    with open(pathlib.Path("./result.jsonl"), "w") as f:
        f.write(result_text)
    
    f.close()